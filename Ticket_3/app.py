# rag-backend/app.py

import os
import time
import psutil
import traceback
from datetime import datetime
from flask import Flask, request, jsonify
from flask_cors import CORS

# Import necessary LangChain components
try:
    from pptx import Presentation
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    from sentence_transformers import SentenceTransformer
    from langchain_community.vectorstores.faiss import FAISS
    from langchain.docstore.document import Document
    # Corrected import for InMemoryDocstore to address deprecation warning
    from langchain_community.docstore.in_memory import InMemoryDocstore
    import numpy as np
    import faiss
    from langchain_cohere import ChatCohere
    from langchain_core.prompts import ChatPromptTemplate
    from pydantic import BaseModel, Field
    from langchain_core.messages import HumanMessage
    from langchain_core.output_parsers import StrOutputParser
    from langchain_community.tools.tavily_search import TavilySearchResults
    from langgraph.graph import END, StateGraph, START
    from typing import List, TypedDict
    # Import specific Cohere error to catch it
    from cohere.errors import TooManyRequestsError
    # Import for JSONDecodeError
    import json
except ImportError as e:
    print(f"Erreur d'importation des dépendances LangChain: {e}")
    print("Veuillez vous assurer que toutes les dépendances sont installées. Exécutez 'pip install -r requirements.txt'")
    exit(1)

# LangSmith and API Configuration
# These are typically set as environment variables in a production environment
# For local testing, you can uncomment and set them directly, but use environment variables for deployment.
# os.environ["LANGCHAIN_TRACING_V2"] = "true"
# os.environ["LANGCHAIN_ENDPOINT"] = "https://api.smith.langchain.com"
# os.environ["LANGCHAIN_API_KEY"] = os.getenv("LANGCHAIN_API_KEY", "YOUR_LANGCHAIN_API_KEY") # Replace with your LangChain API key
os.environ["COHERE_API_KEY"] = os.getenv("COHERE_API_KEY", "skzV7D5L3wScCccmIXgLhMEmyyXy9PsZKeal5AYV") # Replace with your Cohere API key
os.environ['TAVILY_API_KEY'] = os.getenv("TAVILY_API_KEY", "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") # Replace with your Tavily API key

# --- Global Rate Limit Control for Cohere API ---
# Trial key limit is 10 API calls / minute
API_CALL_COOLDOWN_SECONDS = 60 / 9  # Aim for 9 calls per minute to be safe
last_cohere_api_call_time = 0

def enforce_cohere_rate_limit():
    global last_cohere_api_call_time
    current_time = time.time()
    elapsed_time = current_time - last_cohere_api_call_time

    if elapsed_time < API_CALL_COOLDOWN_SECONDS:
        sleep_duration = API_CALL_COOLDOWN_SECONDS - elapsed_time
        print(f"--- Enforcing Cohere rate limit: Sleeping for {sleep_duration:.2f} seconds ---")
        time.sleep(sleep_duration)
    last_cohere_api_call_time = time.time()

# --- 1. UNIFIED TAXONOMY ---
OFFER_TAXONOMY = {
    "cybersecurity": {
        "threat_detection": ["SIEM", "SOC", "threat_hunting", "anomaly_detection"],
        "vulnerability_management": ["penetration_testing", "security_audit", "risk_assessment"],
        "compliance": ["GDPR", "ISO27001", "SOX", "PCI_DSS"],
        "incident_response": ["forensics", "recovery", "investigation", "containment"],
        "identity_management": ["SSO", "MFA", "privileged_access", "identity_governance"]
    },
    "ai_solutions": {
        "machine_learning": ["supervised", "unsupervised", "reinforcement", "deep_learning"],
        "nlp": ["chatbots", "document_processing", "sentiment_analysis", "language_models"],
        "computer_vision": ["object_detection", "image_classification", "facial_recognition"],
        "predictive_analytics": ["forecasting", "anomaly_detection", "pattern_recognition"]
    }
}

# --- 2. PERFORMANCE METRICS ---
class PerformanceMetrics:
    def __init__(self):
        self.start_time = None
        self.end_time = None
        self.memory_start = None

    def start_measurement(self):
        self.start_time = time.time()
        self.memory_start = psutil.Process().memory_info().rss / 1024 / 1024  # MB

    def end_measurement(self):
        self.end_time = time.time()
        memory_end = psutil.Process().memory_info().rss / 1024 / 1024  # MB

        return {
            "latency_seconds": round(self.end_time - self.start_time, 3),
            "memory_used_mb": round(memory_end - self.memory_start, 2),
            "timestamp": datetime.now().isoformat()
        }

# --- 3. COST TRACKING ---
class CostTracker:
    def __init__(self):
        self.cohere_costs = {
            "command-r": {"input": 0.0015, "output": 0.002},  # par 1K tokens
        }
        self.tavily_cost = 0.001  # par recherche
        self.total_cost = 0
        self.cost_breakdown = []

    def add_cohere_cost(self, input_tokens, output_tokens, model="command-r"):
        cost = (input_tokens/1000 * self.cohere_costs[model]["input"] +
                output_tokens/1000 * self.cohere_costs[model]["output"])
        self.total_cost += cost
        self.cost_breakdown.append({
            "service": "cohere",
            "model": model,
            "cost": cost,
            "tokens": {"input": input_tokens, "output": output_tokens}
        })
        return cost

    def add_search_cost(self):
        self.total_cost += self.tavily_cost
        self.cost_breakdown.append({
            "service": "tavily",
            "cost": self.tavily_cost
        })

    def get_total_cost(self):
        return round(self.total_cost, 6)

# --- 4. QUALITY SCORER ---
class QualityScorer:
    def __init__(self):
        self.weights = {
            "relevance": 0.3,
            "completeness": 0.25,
            "coherence": 0.25,
            "technical_accuracy": 0.2
        }

    def calculate_offer_quality(self, generated_offer, question):
        scores = {
            "relevance": self._score_relevance(generated_offer, question),
            "completeness": self._score_completeness(generated_offer),
            "coherence": self._score_coherence(generated_offer),
            "technical_accuracy": self._score_technical_accuracy(generated_offer)
        }

        overall_score = sum(scores[metric] * self.weights[metric] for metric in scores)

        return {
            "overall_score": round(overall_score, 2),
            "detailed_scores": scores,
            "grade": self._get_grade(overall_score)
        }

    def _score_relevance(self, offer, question):
        question_lower = question.lower()
        offer_lower = offer.lower()

        cyber_terms = ["cybersecurity", "cyber", "security", "threat", "vulnerability", "compliance"]
        ai_terms = ["ai", "artificial intelligence", "machine learning", "automation"]

        score = 0
        if any(term in question_lower for term in cyber_terms):
            score += 0.5 if any(term in offer_lower for term in cyber_terms) else 0
        if any(term in question_lower for term in ai_terms):
            score += 0.5 if any(term in offer_lower for term in ai_terms) else 0

        return min(score * 2, 1.0)

    def _score_completeness(self, offer):
        required_sections = ["summary", "solution", "plan", "timeline", "team", "outcomes", "pricing"]
        score = sum(1 for section in required_sections if section.lower() in offer.lower()) / len(required_sections)
        return score

    def _score_coherence(self, offer):
        sentences = offer.split('.')
        if len(sentences) > 3 and len(offer) > 200:
            return 0.8
        elif len(sentences) > 1 and len(offer) > 100:
            return 0.6
        else:
            return 0.4

    def _score_technical_accuracy(self, offer):
        technical_terms = ["implementation", "architecture", "integration", "deployment", "monitoring", "framework", "protocol"]
        score = sum(1 for term in technical_terms if term in offer.lower()) / len(technical_terms)
        return score

    def _get_grade(self, score):
        if score >= 0.9: return "A"
        elif score >= 0.8: return "B"
        elif score >= 0.7: return "C"
        elif score >= 0.6: return "D"
        else: return "F"

# --- 5. TAXONOMY CLASSIFICATION ---
def classify_offer_taxonomy(question):
    question_lower = question.lower()
    classification = {"primary": None, "secondary": [], "confidence": 0}

    for category, subcategories in OFFER_TAXONOMY.items():
        category_matched = False
        if category.replace("_", " ") in question_lower:
            classification["primary"] = category
            classification["confidence"] = 0.8
            category_matched = True

        for subcat_name, subcat_terms in subcategories.items():
            if any(term.lower() in question_lower for term in subcat_terms):
                if not category_matched:
                    classification["primary"] = category
                    classification["confidence"] = 0.6
                if subcat_name not in classification["secondary"]:
                    classification["secondary"].append(subcat_name)

    if classification["primary"] is None:
        if any(term in question_lower for term in ["cybersecurity", "security", "threat", "vulnerability"]):
            classification["primary"] = "cybersecurity"
            classification["confidence"] = 0.4
        elif any(term in question_lower for term in ["ai", "artificial intelligence", "machine learning", "nlp", "computer vision"]):
            classification["primary"] = "ai_solutions"
            classification["confidence"] = 0.4

    return classification

# --- Global RAG Components (initialized once) ---
embedding_model = None
vectorstore = None
retriever = None
llm_router = None
structured_llm_router = None
question_router = None
llm_grader = None
structured_llm_grader = None
retrieval_grader = None
llm_fallback_model = None
llm_chain = None
llm_hallucination_grader = None
structured_llm_hallucination_grader = None
hallucination_grader = None
llm_answer_grader = None
structured_llm_answer_grader = None
answer_grader = None
web_search_tool = None
app_rag_workflow = None

# Flag to track if RAG components are initialized
rag_components_initialized = False

def _initialize_rag_components():
    """
    Initialise tous les composants RAG lourds (modèle d'embedding, index FAISS, LLMs, graphe)
    Cette fonction ne doit être appelée qu'une seule fois au démarrage de l'application Flask.
    """
    global embedding_model, vectorstore, retriever, llm_router, \
           structured_llm_router, question_router, llm_grader, \
           structured_llm_grader, retrieval_grader, llm_fallback_model, \
           llm_chain, llm_hallucination_grader, structured_llm_hallucination_grader, \
           hallucination_grader, llm_answer_grader, structured_llm_answer_grader, \
           answer_grader, web_search_tool, app_rag_workflow, rag_components_initialized

    if rag_components_initialized:
        print("--- Composants RAG déjà initialisés. Saut de la réinitialisation. ---")
        return

    print("--- Initialisation des composants RAG (cela se produit une fois au démarrage du serveur) ---")

    try:
        # Check if API keys are set
        cohere_api_key = os.environ.get("COHERE_API_KEY")
        tavily_api_key = os.environ.get("TAVILY_API_KEY")

        if not cohere_api_key or cohere_api_key == "YOUR_COHERE_API_KEY":
            raise ValueError("COHERE_API_KEY n'est pas définie ou est un placeholder. Veuillez la définir dans vos variables d'environnement ou directement dans app.py.")
        if not tavily_api_key or tavily_api_key == "YOUR_TAVILY_API_KEY":
            raise ValueError("TAVILY_API_KEY n'est pas définie ou est un placeholder. Veuillez la définir dans vos variables d'environnement ou directement dans app.py.")

        # 1. Charger le modèle sentence-transformer
        print("Chargement du modèle d'embedding...")
        embedding_model = SentenceTransformer("BAAI/bge-base-en-v1.5")
        print("Modèle d'embedding chargé.")

        # 2. Définir le chemin PPTX
        # IMPORTANT: Mettez à jour ce chemin vers votre dossier réel contenant les fichiers PPTX
        drive_path = r"C:\Users\legion 5\Desktop\pptxFiles"
        if not os.path.exists(drive_path):
            print(f"AVERTISSEMENT: Le chemin du dossier PPTX n'existe pas: {drive_path}. Le RAG utilisera un vectorstore vide si aucun document n'est trouvé.")
            documents = []
        else:
            # 3. Extracteur PPTX (définition de fonction - déjà globale, mais bon de s'assurer ici)
            def extract_text_from_pptx(file_path):
                prs = Presentation(file_path)
                extracted_text = "\n".join(
                    shape.text.strip()
                    for slide in prs.slides
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text
                )
                return extracted_text

            # 4. Charger les documents
            documents = []
            print(f"Recherche de fichiers PPTX dans : {drive_path}")
            for root, _, files in os.walk(drive_path):
                for fname in files:
                    if fname.lower().endswith(".pptx"):
                        fpath = os.path.join(root, fname)
                        try:
                            text = extract_text_from_pptx(fpath)
                            if text.strip():
                                documents.append(Document(page_content=text, metadata={"source": fpath}))
                            else:
                                print(f"Avertissement : Aucun texte extractible de {fpath}, ignoré.")
                        except Exception as e:
                            print(f"Erreur de lecture de {fpath}: {e}")
                            traceback.print_exc()

        if not documents:
            print("AVERTISSEMENT: Aucun document chargé à partir des fichiers PPTX. Le RAG fonctionnera en mode limité (fallback LLM et recherche web si configurés).")
            vectorstore = None
            retriever = None
        else:
            print(f"Chargé {len(documents)} documents à partir des fichiers PPTX.")

            # 5. Diviser les documents en morceaux
            print("Division des documents en morceaux...")
            splitter = RecursiveCharacterTextSplitter(chunk_size=512, chunk_overlap=50)
            doc_chunks = splitter.split_documents(documents)
            print(f"Documents divisés en {len(doc_chunks)} morceaux.")
            if not doc_chunks:
                print("AVERTISSEMENT: Aucun morceau de document généré. Le vectorstore sera vide.")
                vectorstore = None
                retriever = None
            else:
                # 6. Fonction d'embedding
                class EmbeddingFunction:
                    def __init__(self, model):
                        self.model = model
                    def __call__(self, text):
                        if isinstance(text, str):
                            return self.embed_query(text)
                        else:
                            return self.embed_documents(text)
                    def embed_documents(self, texts):
                        embeddings_tensor = self.model.encode(texts, convert_to_tensor=True)
                        embeddings_np = embeddings_tensor.cpu().detach().numpy()
                        return embeddings_np.astype("float32").tolist()
                    def embed_query(self, text):
                        embeddings_tensor = self.model.encode([text], convert_to_tensor=True)
                        embeddings_np = embeddings_tensor.cpu().detach().numpy()
                        return embeddings_np.astype("float32").tolist()[0]
                embedding_function = EmbeddingFunction(embedding_model)

                # 7. Embed les morceaux de document
                texts_for_embedding = [doc.page_content for doc in doc_chunks]
                print(f"Embedding de {len(texts_for_embedding)} morceaux de document...")
                if not texts_for_embedding:
                    print("AVERTISSEMENT: Aucun texte à embedder. Le vectorstore sera vide.")
                    vectorstore = None
                    retriever = None
                else:
                    embeddings_np = embedding_model.encode(texts_for_embedding, convert_to_tensor=True).cpu().detach().numpy().astype("float32")
                    print("Embedding terminé.")

                    # 8. Construire l'index FAISS
                    print("Construction de l'index FAISS...")
                    if embeddings_np.shape[0] > 0:
                        dimension = embeddings_np.shape[1]
                        index = faiss.IndexFlatL2(dimension)
                        index.add(embeddings_np)
                        print("Index FAISS construit.")
                    else:
                        print("AVERTISSEMENT: Impossible de construire l'index FAISS car aucun embedding n'a été généré.")
                        index = None

                    # 9. Créer le vectorstore FAISS
                    if index is not None:
                        print("Création du vectorstore FAISS...")
                        docstore = InMemoryDocstore({str(i): doc for i, doc in enumerate(doc_chunks)})
                        index_to_docstore_id = {i: str(i) for i in range(len(doc_chunks))}
                        vectorstore = FAISS(
                            embedding_function=embedding_function,
                            index=index,
                            docstore=docstore,
                            index_to_docstore_id=index_to_docstore_id,
                        )
                        print("Vectorstore FAISS créé.")
                        # 10. Construire le retriever
                        retriever = vectorstore.as_retriever()
                        print("Retriever construit.")
                    else:
                        vectorstore = None
                        retriever = None


        # Configuration du LLM et du routeur
        class WebSearch(BaseModel):
            """Utiliser pour les questions NON liées aux offres d'IA, aux fondations d'IA, aux données de pratique Devoteam ou aux sujets de cybersécurité."""
            query: str = Field(description="La requête à utiliser lors de la recherche sur Internet.")

        class VectorStore(BaseModel):
            """Utiliser pour les questions liées à la génération d'offres d'IA, aux fondations d'IA, aux données de pratique Devoteam et aux sujets de cybersécurité. Ceci inclut les questions spécifiques sur les technologies ou services proposés par Devoteam dans ces domaines."""
            query: str = Field(description="La requête à utiliser lors de la recherche dans le vectorstore.")

        preamble_llm_router = """
        Vous êtes un assistant expert en routage de questions. Votre tâche est de diriger la question de l'utilisateur vers l'outil approprié :
        - 'VectorStore' si la question est spécifiquement liée aux offres d'IA de cybersécurité, aux fondations d'IA, aux données internes de Devoteam ou à des sujets de cybersécurité pour lesquels nous devrions avoir des documents internes.
        - 'WebSearch' si la question est de nature générale, nécessite des informations externes, ou n'est manifestement pas couverte par nos documents internes.
        Priorisez 'VectorStore' si la question semble correspondre à notre expertise ou à nos documents.
        """

        llm_router = ChatCohere(model="command-r", temperature=0, cohere_api_key=cohere_api_key)
        structured_llm_router = llm_router.bind_tools(
            tools=[WebSearch, VectorStore],
            preamble=preamble_llm_router
        )
        route_prompt = ChatPromptTemplate.from_messages([("human", "{question}")])
        question_router = route_prompt | structured_llm_router

        # Grader de récupération
        class GradeDocuments(BaseModel):
            binary_score: str = Field(description="Les documents sont pertinents pour la question, 'yes' ou 'no'", enum=["yes", "no"])
        grade_prompt = ChatPromptTemplate.from_messages([
            ("system", """Vous êtes un évaluateur qui évalue la pertinence d'un document récupéré par rapport à une question d'utilisateur.
            Votre tâche est de déterminer si le document contient des informations pertinentes pour répondre à la question de l'utilisateur.
            - Si le document contient des mots-clés, des concepts ou une signification sémantique liés à la question de l'utilisateur, répondez par 'yes'
            - Si le document ne contient pas d'informations pertinentes, répondez par 'no'
            Vous devez répondre UNIQUEMENT par 'yes' ou 'no' - aucun autre texte ou explication."""),
            ("human", "Document: {document}\n\nQuestion: {question}\n\nCe document est-il pertinent pour la question ? Répondez par 'yes' ou 'no' :")
        ])
        llm_grader = ChatCohere(model="command-r", temperature=0, cohere_api_key=cohere_api_key)
        structured_llm_grader = llm_grader.with_structured_output(GradeDocuments)
        retrieval_grader = grade_prompt | structured_llm_grader

        # Fallback LLM
        preamble_llm_fallback = """Vous êtes un assistant pour les tâches de question-réponse. Vous fournissez des réponses générales basées sur vos connaissances générales lorsque des informations spécifiques ne sont pas disponibles. Répondez à la question en vous basant sur vos connaissances. Utilisez trois phrases maximum et soyez concis."""
        llm_fallback_model = ChatCohere(model_name="command-r", temperature=0, cohere_api_key=cohere_api_key).bind(preamble=preamble_llm_fallback)
        llm_chain = ChatPromptTemplate.from_messages([("human", "{question} \nAnswer: ")]) | llm_fallback_model | StrOutputParser()

        # Grader d'hallucination
        class GradeHallucinations(BaseModel):
            """Score binaire pour l'hallucination présente dans la réponse générée."""
            binary_score: str = Field(description="La réponse est basée sur les faits, 'yes' ou 'no'", enum=["yes", "no"])
        hallucination_prompt = ChatPromptTemplate.from_messages([
            ("system", """Vous êtes un évaluateur qui évalue si une génération LLM est basée sur / supportée par un ensemble de faits récupérés.
            Votre tâche est de déterminer si la génération est factuellement basée sur les documents fournis.
            - Si la génération est supportée par les documents, répondez par 'yes'
            - Si la génération contient des informations non trouvées dans les documents, répondez par 'no'
            Vous devez répondre UNIQUEMENT par 'yes' ou 'no' - aucun autre texte ou explication."""),
            ("human", "Documents: {documents}\n\nGénération LLM: {generation}\n\nLa génération est-elle basée sur les documents ? Répondez par 'yes' ou 'no' :")
        ])
        llm_hallucination_grader = ChatCohere(model="command-r", temperature=0, cohere_api_key=cohere_api_key)
        structured_llm_hallucination_grader = llm_hallucination_grader.with_structured_output(GradeHallucinations)
        hallucination_grader = hallucination_prompt | structured_llm_hallucination_grader

        # Grader de réponse
        class GradeAnswer(BaseModel):
            """Score binaire pour évaluer si la réponse répond à la question."""
            binary_score: str = Field(description="La réponse répond à la question, 'yes' ou 'no'", enum=["yes", "no"])
        answer_prompt = ChatPromptTemplate.from_messages([
            ("system", """Vous êtes un évaluateur qui évalue si une réponse répond / résout une question.
            Votre tâche est de déterminer si la réponse répond correctement à la question de l'utilisateur.
            - Si la réponse répond et résout la question, répondez par 'yes'
            - Si la réponse ne répond pas ou ne résout pas la question, répondez par 'no'
            Vous devez répondre UNIQUEMENT par 'yes' ou 'no' - aucun autre texte ou explication."""),
            ("human", "Question: {question}\n\nRéponse: {generation}\n\nLa réponse répond-elle à la question ? Répondez par 'yes' ou 'no' :")
        ])
        llm_answer_grader = ChatCohere(model="command-r", temperature=0, cohere_api_key=cohere_api_key)
        structured_llm_answer_grader = llm_answer_grader.with_structured_output(GradeAnswer)
        answer_grader = answer_prompt | structured_llm_answer_grader

        web_search_tool = TavilySearchResults(api_key=tavily_api_key)

        # État du graphe
        class GraphState(TypedDict):
            """Représente l'état de notre graphe."""
            question: str
            generation: str
            documents: List[Document]
            metrics: dict
            retry_count: int

        # Construire le graphe
        workflow = StateGraph(GraphState, recursion_limit=50)

        # Définir les nœuds
        workflow.add_node("web_search", web_search)
        workflow.add_node("retrieve", retrieve)
        workflow.add_node("grade_documents", grade_documents)
        workflow.add_node("generate", generate)
        workflow.add_node("llm_fallback", llm_fallback)

        # Construire le graphe (arêtes conditionnelles)
        workflow.add_conditional_edges(
            START, route_question,
            {"web_search": "web_search", "vectorstore": "retrieve", "llm_fallback": "llm_fallback"}
        )
        workflow.add_edge("web_search", "generate")
        workflow.add_edge("retrieve", "grade_documents")
        workflow.add_conditional_edges(
            "grade_documents", decide_to_generate,
            {"web_search": "web_search", "generate": "generate", "llm_fallback": "llm_fallback"}
        )
        workflow.add_conditional_edges(
            "generate", grade_generation_v_documents_and_question,
            {"not supported": "generate", "not useful": "web_search", "useful": END, "max_retries_reached": "llm_fallback"}
        )
        workflow.add_edge("llm_fallback", END)

        global app_rag_workflow
        app_rag_workflow = workflow.compile()
        rag_components_initialized = True
        print("--- Composants RAG initialisés avec succès ---")

    except Exception as e:
        print(f"ERREUR CRITIQUE lors de l'initialisation des composants RAG: {e}")
        traceback.print_exc()
        rag_components_initialized = False


# --- Instances globales pour le suivi des métriques (réinitialisées par requête) ---
performance_metrics = PerformanceMetrics()
cost_tracker = CostTracker()

# --- Fonctions de workflow du graphe (ces fonctions utilisent les composants globals) ---
def retrieve(state):
    print("---RÉCUPÉRATION---")
    question = state["question"]
    documents = []
    if retriever is None:
        print("AVERTISSEMENT: Le retriever n'est pas initialisé. Aucun document ne sera récupéré.")
        return {"documents": [], "question": question, "retry_count": state.get("retry_count", 0)}
    try:
        documents = retriever.invoke(question)
        print(f"Récupéré {len(documents)} documents.")
    except Exception as e:
        print(f"Erreur lors de la récupération des documents: {e}")
        traceback.print_exc()
        documents = []
    return {"documents": documents, "question": question, "retry_count": state.get("retry_count", 0)}

def llm_fallback(state):
    print("---Fallback LLM---")
    enforce_cohere_rate_limit() # Enforce rate limit before calling Cohere
    question = state["question"]
    generation = ""
    try:
        input_for_fallback = state.get("generation", question)
        if not input_for_fallback.strip():
             input_for_fallback = question

        generation = llm_chain.invoke({"question": input_for_fallback})
        final_metrics = {
            "performance": performance_metrics.end_measurement(),
            "quality": QualityScorer().calculate_offer_quality(generation, question),
            "cost": cost_tracker.get_total_cost(),
            "taxonomy": classify_offer_taxonomy(question)
        }
        return {"question": question, "generation": generation, "metrics": final_metrics, "retry_count": state.get("retry_count", 0)}
    except TooManyRequestsError as e:
        print(f"Erreur Cohere Rate Limit lors du fallback LLM: {e}")
        return {"question": question, "generation": "Le service LLM est temporairement surchargé (code 429). Veuillez réessayer dans un instant.", "metrics": {}, "retry_count": state.get("retry_count", 0)}
    except Exception as e:
        print(f"Erreur lors du fallback LLM: {e}")
        traceback.print_exc()
        return {"question": question, "generation": "Une erreur est survenue lors de la tentative de réponse générale. Veuillez réessayer.", "metrics": {}, "retry_count": state.get("retry_count", 0)}


def generate(state):
    return generate_with_metrics(state)

def grade_documents(state):
    print("---VÉRIFICATION DE LA PERTINENCE DU DOCUMENT PAR RAPPORT À LA QUESTION---")
    enforce_cohere_rate_limit() # Enforce rate limit before calling Cohere
    question = state["question"]
    documents = state["documents"]
    filtered_docs = []
    if not documents:
        print("---Aucun document à noter, passez à la recherche web ou au fallback.---")
        return {"documents": [], "question": question, "retry_count": state.get("retry_count", 0)}

    for i, d in enumerate(documents):
        try:
            score_obj = retrieval_grader.invoke({"question": question, "document": d.page_content})
            grade = score_obj.binary_score if hasattr(score_obj, 'binary_score') else score_obj.get('binary_score', 'no')
        except TooManyRequestsError as e:
            print(f"Erreur Cohere Rate Limit lors de la notation du document: {e}")
            grade = "no"
            # If rate limited here, it's critical. We'll mark this doc as non-relevant and proceed.
            # The rate limit enforcement should ideally prevent this for subsequent calls.
        except Exception as e:
            print(f"Erreur lors de la notation de la pertinence du document {i}: {e}")
            traceback.print_exc()
            grade = "no"

        if grade == "yes":
            print(f"---NOTE : DOCUMENT PERTINENT ({d.metadata.get('source', 'unknown source')})---")
            filtered_docs.append(d)
        else:
            print(f"---NOTE : DOCUMENT NON PERTINENT ({d.metadata.get('source', 'unknown source')})---")
            continue
    return {"documents": filtered_docs, "question": question, "retry_count": state.get("retry_count", 0)}

def web_search(state):
    print("---RECHERCHE WEB---")
    question = state["question"]
    global cost_tracker
    web_documents = []
    try:
        cost_tracker.add_search_cost() # Tavily has its own rate limits, Cohere limit doesn't apply here directly
        docs = web_search_tool.invoke({"query": question})
        for d in docs:
            web_doc = Document(page_content=d["content"], metadata={"source": d.get("url", "web")})
            web_documents.append(web_doc)
        print(f"Récupéré {len(web_documents)} documents via la recherche web.")
    except Exception as e:
        print(f"Erreur lors de la recherche web: {e}")
        traceback.print_exc()
        web_documents = []
    current_metrics = state.get("metrics", {})
    current_metrics["cost"] = cost_tracker.get_total_cost()
    return {"documents": web_documents, "question": question, "metrics": current_metrics, "retry_count": state.get("retry_count", 0)}

def route_question(state):
    print("---ROUTE DE LA QUESTION---")
    enforce_cohere_rate_limit() # Enforce rate limit before calling Cohere
    question = state["question"]

    state["retry_count"] = state.get("retry_count", 0) + 1
    MAX_ROUTER_RETRIES = 2

    if state["retry_count"] > MAX_ROUTER_RETRIES:
        print(f"---Limite de tentatives de routage ({MAX_ROUTER_RETRIES}) atteinte, bascule vers le fallback LLM.---")
        return "llm_fallback"

    datasource_choice = "llm_fallback"

    try:
        source = question_router.invoke({"question": question})
        tool_calls_from_source = source.response_metadata.get('tool_calls', [])

        if tool_calls_from_source:
            tool_call = tool_calls_from_source[0]
            if isinstance(tool_call, dict) and 'function' in tool_call and 'name' in tool_call['function']:
                datasource_choice = tool_call['function']['name']
                print(f"---Routeur choisi : {datasource_choice}---")
            else:
                print(f"---Structure d'objet d'appel d'outil inattendue: {tool_call}---")
                datasource_choice = "llm_fallback"
        else:
            print("---Aucun appel d'outil détecté par le routeur.---")
            datasource_choice = "llm_fallback"

    except TooManyRequestsError as e:
        print(f"Erreur Cohere Rate Limit lors du routage de la question: {e}")
        return "llm_fallback"
    except Exception as e:
        print(f"Erreur lors du routage de la question: {e}")
        traceback.print_exc()
        datasource_choice = "llm_fallback"

    if datasource_choice == "VectorStore" and retriever is None:
        print("---AVERTISSEMENT: Tentative de routage vers VectorStore mais le retriever n'est pas initialisé. Bascule vers la recherche web.---")
        return "web_search"
    elif datasource_choice == "VectorStore":
        print("---ROUTE DE LA QUESTION VERS RAG---")
        return "vectorstore"
    elif datasource_choice == "WebSearch":
        print("---ROUTE DE LA QUESTION VERS LA RECHERCHE WEB---")
        return "web_search"
    else:
        print(f"---Route de la question vers LLM (Route non gérée ou par défaut)---")
        return "llm_fallback"


def decide_to_generate(state):
    print("---ÉVALUATION DES DOCUMENTS NOTÉS---")
    filtered_documents = state["documents"]
    if not filtered_documents:
        print("---DÉCISION : TOUS LES DOCUMENTS NE SONT PAS PERTINENTS POUR LA QUESTION OU AUCUN DOCUMENT, RECHERCHE WEB OU FALLBACK---\n")
        # Only increment retry_count if this is a true retry due to lack of relevant docs
        if state.get("generation", "").strip() != "": # If we already tried generation and it failed for content reasons
             state["retry_count"] = state.get("retry_count", 0) + 1
        
        MAX_DOC_GRADE_RETRIES = 1

        if state["retry_count"] > MAX_DOC_GRADE_RETRIES:
            print(f"---Limite de tentatives de récupération/recherche web ({MAX_DOC_GRADE_RETRIES}) atteinte, bascule vers le fallback LLM.---")
            return "llm_fallback"
        return "web_search"
    else:
        print("---DÉCISION : GÉNÉRER---")
        # Reset retry_count when we decide to generate based on new relevant documents
        state["retry_count"] = 0 
        return "generate"

def grade_generation_v_documents_and_question(state):
    print("---VÉRIFICATION DES HALLUCINATIONS ET DE LA PERTINENCE DE LA RÉPONSE---")
    enforce_cohere_rate_limit() # Enforce rate limit before calling Cohere
    question = state["question"]
    documents = state["documents"]
    generation = state["generation"]

    state["retry_count"] = state.get("retry_count", 0) + 1
    MAX_GENERATION_RETRIES = 2

    hallucination_grade = "no"
    if not generation or "temporairement surchargé" in generation or "erreur est survenue" in generation:
        print("---Génération vide ou en erreur, pas de notation d'hallucination/réponse.---")
        if state["retry_count"] >= MAX_GENERATION_RETRIES:
            print("---Limite de tentatives de génération atteinte (génération en erreur), bascule vers le fallback LLM.---")
            return "max_retries_reached"
        return "not useful"

    if documents:
        try:
            score_obj = hallucination_grader.invoke({"documents": documents, "generation": generation})
            hallucination_grade = score_obj.binary_score if hasattr(score_obj, 'binary_score') else score_obj.get('binary_score', 'no')
        except TooManyRequestsError as e:
            print(f"Erreur Cohere Rate Limit lors de la notation des hallucinations: {e}")
            if state["retry_count"] >= MAX_GENERATION_RETRIES:
                print("---Limite de tentatives de génération atteinte en raison de Rate Limit, bascule vers le fallback LLM.---")
                return "max_retries_reached"
            return "not supported"
        except json.decoder.JSONDecodeError as e:
            print(f"Erreur JSONDecodeError lors de la notation des hallucinations: {e}")
            print("Cela indique souvent une réponse non-JSON de l'API Cohere (ex: surcharge).")
            if state["retry_count"] >= MAX_GENERATION_RETRIES:
                print("---Limite de tentatives de génération atteinte (JSONDecodeError), bascule vers le fallback LLM.---")
                return "max_retries_reached"
            return "not supported"
        except Exception as e:
            print(f"Erreur lors de la notation des hallucinations: {e}")
            traceback.print_exc()
            hallucination_grade = "no"

    if hallucination_grade == "yes":
        print("---DÉCISION : LA GÉNÉRATION EST BASÉE SUR LES DOCUMENTS---")
        answer_grade = "no"
        try:
            score_obj = answer_grader.invoke({"question": question, "generation": generation})
            answer_grade = score_obj.binary_score if hasattr(score_obj, 'binary_score') else score_obj.get('binary_score', 'no')
        except TooManyRequestsError as e:
            print(f"Erreur Cohere Rate Limit lors de la notation de la réponse: {e}")
            if state["retry_count"] >= MAX_GENERATION_RETRIES:
                print("---Limite de tentatives de génération atteinte en raison de Rate Limit, bascule vers le fallback LLM.---")
                return "max_retries_reached"
            return "not useful"
        except json.decoder.JSONDecodeError as e:
            print(f"Erreur JSONDecodeError lors de la notation de la réponse: {e}")
            print("Cela indique souvent une réponse non-JSON de l'API Cohere (ex: surcharge).")
            if state["retry_count"] >= MAX_GENERATION_RETRIES:
                print("---Limite de tentatives de génération atteinte (JSONDecodeError), bascule vers le fallback LLM.---")
                return "max_retries_reached"
            return "not useful"
        except Exception as e:
            print(f"Erreur lors de la notation de la réponse: {e}")
            traceback.print_exc()
            answer_grade = "no"

        if answer_grade == "yes":
            print("---DÉCISION : LA GÉNÉRATION RÉPOND À LA QUESTION. FIN DU FLUX.---\n")
            return "useful"
        else:
            print("---DÉCISION : LA GÉNÉRATION NE RÉPOND PAS À LA QUESTION.---\n")
            if state["retry_count"] >= MAX_GENERATION_RETRIES:
                print("---Limite de tentatives de génération atteinte, bascule vers le fallback LLM.---")
                return "max_retries_reached"
            return "not useful"
    else:
        print("---DÉCISION : LA GÉNÉRATION N'EST PAS BASÉE SUR LES DOCUMENTS (hallucination ou données insuffisantes).---\n")
        if state["retry_count"] >= MAX_GENERATION_RETRIES:
            print("---Limite de tentatives de génération atteinte, bascule vers le fallback LLM.---")
            return "max_retries_reached"
        return "not supported"

def generate_with_metrics(state):
    global performance_metrics, cost_tracker

    preamble_generate = """Vous êtes un assistant pour la génération d'offres d'IA de cybersécurité. Utilisez les éléments de contexte récupérés suivants pour générer une offre de cybersécurité complète qui comprend :

    1. RÉSUMÉ EXÉCUTIF : Bref aperçu de la solution de cybersécurité alimentée par l'IA
    2. SOLUTION TECHNIQUE : Approche détaillée utilisant l'IA pour la cybersécurité
    3. PLAN DE MISE EN ŒUVRE : Stratégie de déploiement étape par étape
    4. CALENDRIER : Phases du projet et jalons
    5. COMPOSITION DE L'ÉQUIPE : Expertise et rôles requis
    6. RÉSULTATS ATTENDUS : Bénéfices mesurables et ROI
    7. STRUCTURE DE PRIX : Répartition des coûts et investissement

    Générez une offre professionnelle et détaillée en français. Si vous ne connaissez pas les détails spécifiques, utilisez le contexte fourni et les meilleures pratiques générales en matière d'IA de cybersécurité.
    Assurez-vous que toutes les sections demandées sont présentes et clairement numérotées et titrées en gras.
    """

    llm_generate = ChatCohere(model_name="command-r", temperature=0, cohere_api_key=os.environ["COHERE_API_KEY"])

    documents_content = [doc.page_content for doc in state["documents"]] if state["documents"] else []
    documents_str = "\n".join(documents_content)

    rag_chain = ChatPromptTemplate.from_messages([
        HumanMessage(
            f"Question: {state['question']}\nContext: {documents_str}\n\nGenerate a comprehensive cybersecurity AI offer, ensuring all sections from the preamble are included:",
            additional_kwargs={"documents": state["documents"]},
        )
    ]) | llm_generate | StrOutputParser()

    performance_metrics.start_measurement()

    print("---GÉNÉRATION AMÉLIORÉE AVEC DES MÉTRIQUES---")
    question = state["question"]
    documents = state["documents"]

    generation = ""
    try:
        enforce_cohere_rate_limit() # Enforce rate limit before calling Cohere
        generation = rag_chain.invoke({"documents": documents, "question": question})
    except TooManyRequestsError as e:
        print(f"Erreur Cohere Rate Limit lors de la génération de l'offre: {e}")
        generation = "Le service de génération d'offre est temporairement surchargé. Veuillez réessayer dans un instant."
    except Exception as e:
        print(f"Erreur lors de la génération de l'offre: {e}")
        traceback.print_exc()
        generation = "Une erreur est survenue lors de la génération de l'offre. Veuillez réessayer."

    perf_metrics = performance_metrics.end_measurement()

    estimated_input_tokens = len(question.split())
    for doc in documents:
        estimated_input_tokens += len(doc.page_content.split())
    estimated_output_tokens = len(generation.split()) * 1.3

    if generation and "temporairement surchargé" not in generation and "erreur est survenue" not in generation:
        cost_tracker.add_cohere_cost(estimated_input_tokens, estimated_output_tokens)

    quality_scorer_instance = QualityScorer()
    quality_score = quality_scorer_instance.calculate_offer_quality(generation, question)

    taxonomy_class = classify_offer_taxonomy(question)

    combined_metrics = state.get("metrics", {})
    combined_metrics["performance"] = perf_metrics
    combined_metrics["quality"] = quality_score
    combined_metrics["cost"] = cost_tracker.get_total_cost()
    combined_metrics["taxonomy"] = taxonomy_class

    print(f"Latence : {perf_metrics['latency_seconds']}s")
    print(f"Coût total : ${cost_tracker.get_total_cost()}")
    print(f"Score de qualité : {quality_score['overall_score']}/1.0 (Note : {quality_score['grade']})")
    print(f"Taxonomie : {taxonomy_class.get('primary')} - {taxonomy_class.get('secondary')}\n")

    return {
        "question": question,
        "documents": documents,
        "generation": generation,
        "metrics": combined_metrics,
        "retry_count": state.get("retry_count", 0)
    }


# --- Configuration de l'application Flask ---
app = Flask(__name__)
CORS(app)

with app.app_context():
    _initialize_rag_components()

@app.route('/ask-rag', methods=['POST'])
def ask_rag():
    user_question = request.json.get('question')
    if not user_question:
        return jsonify({"error": "Aucune question fournie"}), 400

    global rag_components_initialized
    if not rag_components_initialized:
        print("ERREUR: Le service RAG n'est pas entièrement initialisé. Refus de la requête.")
        return jsonify({"error": "Le service RAG n'est pas entièrement initialisé. Veuillez vérifier les logs du serveur pour plus de détails."}), 503

    try:
        global performance_metrics, cost_tracker
        performance_metrics = PerformanceMetrics()
        cost_tracker = CostTracker()

        inputs = {"question": user_question, "retry_count": 0, "documents": [], "generation": "", "metrics": {}}
        # This will hold the *last* state yielded by the graph.
        # If the graph reaches END, this will be the final state.
        last_known_state = None

        try:
            for output in app_rag_workflow.stream(inputs):
                # LangGraph's .stream() yields a dictionary.
                # 'output' contains the state *after* the current node.
                # Update last_known_state with the most recent full state.
                # If END is reached, output[END] holds the final state.
                if END in output:
                    last_known_state = output[END]
                    break # Exit the loop as the graph has finished
                else:
                    # If not END, it's an intermediate node's output.
                    # We need to extract the relevant state from this output.
                    # LangGraph's stream often yields {node_name: {updated_state_dict}}.
                    # We'll just take the first value as the current state.
                    if output: # Ensure output is not empty
                        last_known_state = list(output.values())[0]

        except TooManyRequestsError as e:
            print(f"Erreur Cohere Rate Limit au niveau du workflow : {e}")
            return jsonify({"error": "Le service LLM est temporairement surchargé (code 429). Veuillez réessayer dans un instant."}), 429
        except json.decoder.JSONDecodeError as e:
            print(f"Erreur JSONDecodeError au niveau du workflow: {e}")
            print("Cela indique souvent une réponse non-JSON de l'API Cohere (ex: surcharge).")
            return jsonify({"error": "Une erreur de communication avec le service LLM est survenue (format de réponse inattendu). Veuillez réessayer."}), 500
        except Exception as e:
            print(f"Erreur inattendue au niveau du workflow: {e}")
            traceback.print_exc()
            return jsonify({"error": f"Une erreur inattendue est survenue lors de l'exécution du workflow: {str(e)}"}), 500

        # After the loop, use last_known_state
        if last_known_state:
            metrics_to_return = last_known_state.get("metrics", {})
            return jsonify({
                "generation": last_known_state.get("generation", "Aucune génération disponible."),
                "metrics": metrics_to_return
            })
        else:
            print("AVERTISSEMENT: Le workflow RAG n'a pas produit de sortie finale. Cela peut indiquer un problème dans le graphe ou une erreur non capturée.")
            return jsonify({"error": "Le workflow RAG n'a pas produit de sortie finale après toutes les tentatives. Veuillez vérifier les logs du serveur."}), 500

    except Exception as e:
        print(f"Erreur lors du traitement RAG global : {e}")
        traceback.print_exc()
        return jsonify({"error": f"Une erreur interne est survenue: {str(e)}"}), 500

if __name__ == '__main__':
    app.run(port=5000, debug=True)